from pptx import Presentation

from services.file.app.parser.base_parser import BaseParser


class PptxParser(BaseParser):


    async def extract_text(self, file_path: str) -> str:
        presentation = Presentation(file_path)
        text = []

        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text.append(shape.text)
        return "\n".join(text)